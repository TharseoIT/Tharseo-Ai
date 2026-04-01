from pathlib import Path
from typing import Literal

FileType = Literal["pdf", "docx", "pptx", "txt"]

ALLOWED_EXTENSIONS: dict[str, FileType] = {
    ".pdf":  "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".txt":  "txt",
}


def extract_text(file_path: Path, file_type: FileType) -> str:
    """Extract raw text from a document file. Returns a single string."""
    if file_type == "txt":
        return file_path.read_text(encoding="utf-8", errors="replace")

    if file_type == "pdf":
        from pypdf import PdfReader
        reader = PdfReader(str(file_path))
        pages = [page.extract_text() for page in reader.pages if page.extract_text()]
        return "\n\n".join(pages)

    if file_type == "docx":
        from docx import Document as DocxDoc
        doc = DocxDoc(str(file_path))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if file_type == "pptx":
        from pptx import Presentation
        prs = Presentation(str(file_path))
        slides = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            texts.append(line)
            if texts:
                slides.append("\n".join(texts))
        return "\n\n".join(slides)

    raise ValueError(f"Unsupported file type: {file_type}")
